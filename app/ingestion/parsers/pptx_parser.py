from pptx import Presentation
from pathlib import Path

from app.models.models import DocumentSection

def parse_pptx(file_path: str) -> list[DocumentSection]:

    prs = Presentation(file_path)

    sections = []

    for slide_num, slide in enumerate(prs.slides, start=1):

        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_text.append(text)

        if not slide_text:
            continue

        sections.append(
            DocumentSection(
                text="\n".join(slide_text),
                source_file=Path(file_path).name,
                document_type="pptx",
                metadata={
                    "slide": slide_num
                }
            )
        )

    return sections

from openpyxl import load_workbook
from pathlib import Path
